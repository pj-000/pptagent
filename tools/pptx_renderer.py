import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from models.schemas import PresentationPlan, SlideSpec, TextElement
import config


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "circle": MSO_SHAPE.OVAL,
    "line": MSO_SHAPE.RECTANGLE,  # 用极窄矩形模拟线条
}


class PPTXRenderer:

    def render(self, plan: PresentationPlan, filename: str = "output.pptx") -> str:
        prs = Presentation()
        prs.slide_width = Inches(plan.slide_width)
        prs.slide_height = Inches(plan.slide_height)

        for slide_spec in plan.slides:
            self._render_slide(prs, slide_spec, plan)

        os.makedirs(config.OUTPUT_DIR, exist_ok=True)
        output_path = os.path.join(config.OUTPUT_DIR, filename)
        prs.save(output_path)
        print(f"[Renderer] 已保存：{output_path}（共 {len(plan.slides)} 页）")
        return os.path.abspath(output_path)

    def _render_slide(self, prs: Presentation, spec: SlideSpec, plan: PresentationPlan):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 背景色
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(spec.background_color)

        # 先渲染形状（在底层），再渲染文字（在上层）
        for elem in spec.elements:
            if elem.type == "shape":
                self._add_shape(slide, elem)

        for elem in spec.elements:
            if elem.type == "image_placeholder":
                self._add_image(slide, elem)
            elif elem.type != "shape":
                self._add_text_box(slide, elem, plan.font_family)

        if spec.speaker_notes:
            slide.notes_slide.notes_text_frame.text = spec.speaker_notes

    def _add_text_box(self, slide, elem: TextElement, font_family: str):
        txBox = slide.shapes.add_textbox(
            Inches(elem.x), Inches(elem.y),
            Inches(elem.width), Inches(elem.height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = elem.content.split("\n") if elem.content else [""]
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = ALIGN_MAP.get(elem.align, PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = line

            font = run.font
            font.name = font_family
            font.size = Pt(elem.font_size)
            font.bold = elem.bold
            font.color.rgb = hex_to_rgb(elem.color)

    def _add_shape(self, slide, elem: TextElement):
        """渲染 rect / circle / line 形状。"""
        shape_type = SHAPE_MAP.get(elem.shape_type, MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(
            shape_type,
            Inches(elem.x), Inches(elem.y),
            Inches(elem.width), Inches(elem.height)
        )

        # 填充色
        if elem.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(elem.fill_color)

        # 无边框
        shape.line.fill.background()

        # 圆角（仅 rect）
        if elem.shape_type == "rect" and elem.corner_radius:
            try:
                from pptx.oxml.ns import qn
                sp = shape._element
                prstGeom = sp.find(qn("a:prstGeom"))
                if prstGeom is not None:
                    avLst = prstGeom.find(qn("a:avLst"))
                    if avLst is None:
                        from lxml import etree
                        avLst = etree.SubElement(prstGeom, qn("a:avLst"))
            except Exception:
                pass  # 圆角是增强功能，失败不影响主流程

        # 如果形状有文字内容（如 callout 内嵌在 shape 里），添加文字
        if elem.content and elem.content.strip():
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = ALIGN_MAP.get(elem.align, PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = elem.content
            run.font.size = Pt(elem.font_size)
            run.font.bold = elem.bold
            run.font.color.rgb = hex_to_rgb(elem.color)

    def _add_image(self, slide, elem: TextElement):
        left = Inches(elem.x)
        top = Inches(elem.y)
        width = Inches(elem.width)
        height = Inches(elem.height)

        if elem.local_image_path and os.path.isfile(elem.local_image_path):
            try:
                slide.shapes.add_picture(elem.local_image_path, left, top, width, height)
                return
            except Exception as e:
                print(f"[Renderer] 图片插入失败，使用占位: {e}")

        # 灰色占位矩形
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        desc = elem.content or elem.unsplash_query or "图片"
        run.text = f"[图片] {desc}"
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

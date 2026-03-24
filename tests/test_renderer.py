import pytest
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from models.schemas import (
    PresentationPlan, SlideSpec, SlideLayout, TextElement
)
from tools.pptx_renderer import PPTXRenderer, hex_to_rgb
from pptx.dml.color import RGBColor


@pytest.fixture
def renderer():
    return PPTXRenderer()


@pytest.fixture
def simple_plan():
    """最简单的三页 PPT，用于基础渲染测试"""
    return PresentationPlan(
        title="测试 PPT",
        topic="测试",
        slides=[
            SlideSpec(
                slide_index=0,
                layout=SlideLayout.COVER,
                topic="封面",
                elements=[
                    TextElement(
                        content="测试标题",
                        x=1.0, y=2.5, width=11.333, height=1.5,
                        font_size=40, bold=True, color="#1F3864", align="center"
                    )
                ]
            ),
            SlideSpec(
                slide_index=1,
                layout=SlideLayout.CONTENT,
                topic="正文",
                elements=[
                    TextElement(
                        content="正文标题",
                        x=0.5, y=0.3, width=12.333, height=0.9,
                        font_size=32, bold=True, color="#1F3864", align="left"
                    ),
                    TextElement(
                        content="• 要点一\n• 要点二\n• 要点三",
                        x=0.5, y=1.5, width=12.333, height=5.5,
                        font_size=18, color="#333333", align="left"
                    )
                ]
            ),
            SlideSpec(
                slide_index=2,
                layout=SlideLayout.CLOSING,
                topic="结尾",
                elements=[
                    TextElement(
                        content="感谢聆听",
                        x=1.0, y=2.5, width=11.333, height=1.2,
                        font_size=44, bold=True, color="#1F3864", align="center"
                    )
                ]
            )
        ]
    )


def test_hex_to_rgb():
    """测试颜色转换函数"""
    color = hex_to_rgb("#1F3864")
    assert color == RGBColor(0x1F, 0x38, 0x64)

    color2 = hex_to_rgb("#FFFFFF")
    assert color2 == RGBColor(255, 255, 255)


def test_render_creates_file(renderer, simple_plan, tmp_path):
    """测试渲染后文件确实存在"""
    import config
    original_output = config.OUTPUT_DIR
    config.OUTPUT_DIR = str(tmp_path)

    try:
        output_path = renderer.render(simple_plan, "test_output.pptx")
        assert os.path.exists(output_path)
        assert output_path.endswith(".pptx")
    finally:
        config.OUTPUT_DIR = original_output


def test_render_slide_count(renderer, simple_plan, tmp_path):
    """测试渲染后幻灯片数量正确"""
    from pptx import Presentation
    import config
    original_output = config.OUTPUT_DIR
    config.OUTPUT_DIR = str(tmp_path)

    try:
        output_path = renderer.render(simple_plan, "test_count.pptx")
        prs = Presentation(output_path)
        assert len(prs.slides) == len(simple_plan.slides)
    finally:
        config.OUTPUT_DIR = original_output


def test_render_all_layouts(renderer, tmp_path):
    """测试所有布局类型都能正常渲染"""
    from tools.templates import apply_template
    import config
    original_output = config.OUTPUT_DIR
    config.OUTPUT_DIR = str(tmp_path)

    try:
        test_slides = [
            apply_template(SlideLayout.COVER, {"title": "封面", "subtitle": "副标题"}, 0),
            apply_template(SlideLayout.TOC, {"title": "目录", "body": ["第一章", "第二章"]}, 1),
            apply_template(SlideLayout.CONTENT, {"title": "正文", "body": "• 要点一\n• 要点二"}, 2),
            apply_template(SlideLayout.TWO_COLUMN, {"title": "对比", "left": "左侧", "right": "右侧"}, 3),
            apply_template(SlideLayout.CLOSING, {"title": "感谢", "subtitle": "提问环节"}, 4),
        ]

        plan = PresentationPlan(title="布局测试", topic="测试", slides=test_slides)
        output_path = renderer.render(plan, "test_layouts.pptx")
        assert os.path.exists(output_path)
    finally:
        config.OUTPUT_DIR = original_output


def test_render_file_size(renderer, simple_plan, tmp_path):
    """测试生成的文件大小合理（大于 10KB）"""
    import config
    original_output = config.OUTPUT_DIR
    config.OUTPUT_DIR = str(tmp_path)

    try:
        output_path = renderer.render(simple_plan, "test_size.pptx")
        file_size = os.path.getsize(output_path)
        assert file_size > 10 * 1024, f"文件太小：{file_size} bytes"
    finally:
        config.OUTPUT_DIR = original_output


def test_multiline_text(renderer, tmp_path):
    """测试多行文本能正常渲染"""
    import config
    original_output = config.OUTPUT_DIR
    config.OUTPUT_DIR = str(tmp_path)

    try:
        multiline_content = "第一行\n第二行\n第三行\n• 要点一\n• 要点二\n• 要点三"
        plan = PresentationPlan(
            title="多行测试", topic="测试",
            slides=[
                SlideSpec(
                    slide_index=0, layout=SlideLayout.CONTENT, topic="多行",
                    elements=[
                        TextElement(content="标题", x=0.5, y=0.3, width=12.0, height=0.9,
                                    font_size=32, bold=True, color="#000000", align="left"),
                        TextElement(content=multiline_content, x=0.5, y=1.5, width=12.0, height=5.5,
                                    font_size=18, color="#333333", align="left"),
                    ]
                ),
                SlideSpec(
                    slide_index=1, layout=SlideLayout.CLOSING, topic="结尾",
                    elements=[
                        TextElement(content="结束", x=1.0, y=2.5, width=11.0, height=1.2,
                                    font_size=40, bold=True, color="#000000", align="center")
                    ]
                )
            ]
        )
        output_path = renderer.render(plan, "test_multiline.pptx")
        assert os.path.exists(output_path)
    finally:
        config.OUTPUT_DIR = original_output


def test_speaker_notes(renderer, tmp_path):
    """测试演讲者备注能正常写入"""
    from pptx import Presentation
    import config
    original_output = config.OUTPUT_DIR
    config.OUTPUT_DIR = str(tmp_path)

    try:
        plan = PresentationPlan(
            title="备注测试", topic="测试",
            slides=[
                SlideSpec(
                    slide_index=0, layout=SlideLayout.CONTENT, topic="有备注的页",
                    speaker_notes="这是演讲者备注",
                    elements=[
                        TextElement(content="标题", x=0.5, y=0.3, width=12.0, height=0.9,
                                    font_size=32, bold=True, color="#000000", align="left"),
                    ]
                ),
                SlideSpec(
                    slide_index=1, layout=SlideLayout.CLOSING, topic="结尾",
                    elements=[
                        TextElement(content="结束", x=1.0, y=2.5, width=11.0, height=1.2,
                                    font_size=40, bold=True, color="#000000", align="center")
                    ]
                )
            ]
        )
        output_path = renderer.render(plan, "test_notes.pptx")
        prs = Presentation(output_path)
        notes_text = prs.slides[0].notes_slide.notes_text_frame.text
        assert "演讲者备注" in notes_text
    finally:
        config.OUTPUT_DIR = original_output
